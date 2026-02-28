import os
import logging
import tempfile
from typing import List, Dict, Any, Optional
from pathlib import Path
import uuid
from datetime import datetime
import asyncio
import numpy as np

# Document processing libraries
from langchain_text_splitters import RecursiveCharacterTextSplitter

# File format libraries
import PyPDF2
from pptx import Presentation
import docx
from pdf2image import convert_from_path

# Google AI libraries
import google.generativeai as genai

logger = logging.getLogger(__name__)

import os
import config

# Configure Google AI
gemini_key = os.getenv("GEMINI_API_KEY")
if gemini_key:
    genai.configure(api_key=gemini_key)
else:
    genai.configure(api_key="AIzaSyAWKOQzKrlYcVy-uFxdmcK5QWlt1LU-gaI")  # fallback

# Google AI Models
VISION_MODEL = "gemini-1.5-flash"
EMBED_MODEL = "models/embedding-001"

class DocumentProcessor:
    def __init__(self):
        # Initialize text splitter with optimal configuration for detailed analysis
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=300,      # Even smaller chunks for better semantic granularity
            chunk_overlap=50,    # Reduced overlap for more distinct chunks
            separators=["\n\n", "\n", ". ", "; ", ", ", " ", ""]  # More granular separators
        )
        
    def extract_digital_pages(self, file_path: str):
        """Extract text from digital PDF pages"""
        try:
            logger.info(f"Attempting digital text extraction from: {file_path}")
            reader = PyPDF2.PdfReader(file_path)
            logger.info(f"PDF has {len(reader.pages)} pages")
            
            pages = []
            total_text_length = 0

            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    pages.append((i+1, text))
                    total_text_length += len(text.strip())
                    logger.info(f"Page {i+1}: {len(text.strip())} characters extracted")

            logger.info(f"Digital extraction completed: {len(pages)} pages with text, total {total_text_length} characters")
            return pages
        except Exception as e:
            logger.error(f"Error extracting digital pages: {str(e)}")
            return []

    def extract_scanned_pages(self, file_path: str):
        """Extract text from scanned PDF pages using Gemini Vision with detailed analysis"""
        try:
            images = convert_from_path(file_path, dpi=300)
            model = genai.GenerativeModel(VISION_MODEL)
            pages = []

            for i, img in enumerate(images):
                logger.info(f"Analyzing page {i+1} with Gemini Vision")
                prompt = """Extract all readable text from this page. Also describe any images, charts, diagrams, or visual elements present. 
                Provide a detailed analysis including:
                1. All readable text content
                2. Description of any images or visual elements
                3. Any charts, graphs, or diagrams with their content
                4. Overall context and meaning of the page
                Do NOT summarize - provide comprehensive details."""
                response = model.generate_content([prompt, img])
                pages.append((i+1, response.text))

            return pages
        except Exception as e:
            logger.error(f"Error extracting scanned pages: {str(e)}")
            return []

    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file using enhanced method"""
        try:
            logger.info(f"Processing PDF file: {file_path}")
            
            # Try digital text extraction first
            pages = self.extract_digital_pages(file_path)
            logger.info(f"Found {len(pages)} pages with digital extraction")
            
            if len(pages) > 0:
                # Digital PDF - extract text with page numbers
                text = ""
                for page_num, page_text in pages:
                    text += f"Page {page_num}:\n{page_text}\n\n"
                logger.info(f"Successfully extracted text from {len(pages)} pages")
                return text
            else:
                # Scanned PDF - use Gemini Vision OCR
                logger.info("No digital text found, switching to Gemini Vision OCR...")
                pages = self.extract_scanned_pages(file_path)
                logger.info(f"OCR processed {len(pages)} pages")
                
                text = ""
                for page_num, page_text in pages:
                    text += f"Page {page_num}:\n{page_text}\n\n"
                return text
            
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            return ""
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint file with image analysis"""
        try:
            text = ""
            presentation = Presentation(file_path)
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = ""
                image_descriptions = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                    
                    # Extract and analyze images
                    if shape.shape_type == 13:  # Picture shape type
                        try:
                            # Extract image from slide
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Analyze image with Gemini Vision
                            model = genai.GenerativeModel(VISION_MODEL)
                            prompt = """Analyze this image and provide:
                            1. Detailed description of what you see
                            2. Any text visible in the image
                            3. Charts, graphs, or diagrams with their meaning
                            4. Overall context and relevance to the presentation"""
                            
                            response = model.generate_content([prompt, image_bytes])
                            image_descriptions.append(f"Image analysis: {response.text}")
                            
                        except Exception as img_error:
                            logger.warning(f"Could not analyze image on slide {slide_num}: {str(img_error)}")
                
                # Combine text and image descriptions
                if slide_text.strip() or image_descriptions:
                    slide_content = f"Slide {slide_num}:\n"
                    if slide_text.strip():
                        slide_content += f"Text content:\n{slide_text}\n"
                    if image_descriptions:
                        slide_content += f"Image descriptions:\n" + "\n".join(image_descriptions) + "\n"
                    slide_content += "\n"
                    text += slide_content
            
            return text
            
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            return ""
    
    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from Word document with image analysis"""
        try:
            text = ""
            doc = docx.Document(file_path)
            
            # Extract text from paragraphs
            for para_num, paragraph in enumerate(doc.paragraphs, 1):
                if paragraph.text.strip():
                    text += f"Paragraph {para_num}: {paragraph.text}\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text += f"Table cell: {cell.text}\n"
            
            # Extract and analyze images
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        # Get image data
                        image_part = rel.target_part
                        image_bytes = image_part.blob
                        
                        # Analyze image with Gemini Vision
                        model = genai.GenerativeModel(VISION_MODEL)
                        prompt = """Analyze this image and provide:
                        1. Detailed description of what you see
                        2. Any text visible in the image
                        3. Charts, graphs, or diagrams with their meaning
                        4. Overall context and relevance to the document"""
                        
                        response = model.generate_content([prompt, image_bytes])
                        text += f"Image analysis: {response.text}\n\n"
                        
                    except Exception as img_error:
                        logger.warning(f"Could not analyze image in DOCX: {str(img_error)}")
            
            return text
            
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            return ""
    
    def extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                return f"Text file content:\n{content}"
                
        except Exception as e:
            logger.error(f"Error extracting text from TXT: {str(e)}")
            return ""
    
    def extract_text(self, file_path: str, file_type: str) -> str:
        """Extract text based on file type"""
        logger.info(f"Extracting text from {file_type} file: {file_path}")
        
        if file_type.lower() == 'pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_type.lower() in ['pptx', 'ppt']:
            return self.extract_text_from_pptx(file_path)
        elif file_type.lower() in ['docx', 'doc']:
            return self.extract_text_from_docx(file_path)
        elif file_type.lower() in ['txt', 'text']:
            return self.extract_text_from_txt(file_path)
        else:
            logger.error(f"Unsupported file type: {file_type}")
            return ""
    
    def get_embedding(self, text: str):
        """Get embedding using Google Gemini"""
        try:
            response = genai.embed_content(
                model="models/gemini-embedding-001",
                content=text
            )
            embedding = np.array(response["embedding"]).astype("float32")
            logger.info(f"Generated embedding for text: '{text[:50]}...' - length: {len(embedding)}")
            return embedding
        except Exception as e:
            logger.error(f"Error getting embedding: {str(e)}")
            return None

    async def process_document(self, file, user_id: str, document_id: str, 
                        filename: str) -> Dict[str, Any]:
        """Process uploaded document and store embeddings"""
        start_time = datetime.utcnow()
        
        try:
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{filename}") as temp_file:
                temp_file.write(file.file.read())
                temp_file_path = temp_file.name
            
            # Determine file type
            file_extension = Path(filename).suffix.lower().lstrip('.')
            
            # Extract text
            text = self.extract_text(temp_file_path, file_extension)
            
            if not text.strip():
                raise Exception("Could not extract text from document")
            
            # Split text into chunks
            chunks = self.text_splitter.create_documents([text])
            
            # Convert to our format and generate embeddings
            chunk_data = []
            embeddings = []
            
            for i, chunk in enumerate(chunks):
                chunk_text = chunk.page_content
                embedding = self.get_embedding(chunk_text)
                
                if embedding is not None:
                    chunk_data.append({
                        "content": chunk_text,
                        "metadata": {
                            "chunk_index": i,
                            "source": filename
                        }
                    })
                    embeddings.append(embedding.tolist())
            
            # Store in vector database
            vector_db = await get_vector_db()
            success = await vector_db.store_embeddings(
                user_id=user_id,
                document_id=document_id,
                filename=filename,
                chunks=chunk_data,
                embeddings=embeddings
            )
            
            # Clean up temp file
            os.unlink(temp_file_path)
            
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            
            return {
                "success": success,
                "chunks_processed": len(chunk_data),
                "processing_time": processing_time,
                "document_id": document_id,
                "filename": filename
            }
            
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            # Clean up temp file on error
            if 'temp_file_path' in locals():
                try:
                    os.unlink(temp_file_path)
                except:
                    pass
            raise e
    
    async def get_user_documents(self, user_id: str) -> List[Dict[str, Any]]:
        """Get all documents for a user"""
        try:
            vector_db = await get_vector_db()
            documents = await vector_db.get_user_documents(user_id)
            return documents
        except Exception as e:
            logger.error(f"Error getting user documents: {str(e)}")
            return []
    
    async def delete_document(self, document_id: str, user_id: str) -> bool:
        """Delete a document and its embeddings"""
        try:
            vector_db = await get_vector_db()
            success = await vector_db.delete_document(user_id, document_id)
            return success
        except Exception as e:
            logger.error(f"Error deleting document: {str(e)}")
            return False

# Import here to avoid circular import
from database import get_vector_db
